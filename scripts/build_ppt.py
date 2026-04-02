"""
Build presentation: Crypto Trading Intelligence Terminal
NMIMS Innovathon 2026 — Challenge 2
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import sys

# ── Palette ─────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0d, 0x11, 0x17)   # #0d1117  github dark
PANEL    = RGBColor(0x07, 0x0b, 0x14)   # #070b14
BORDER   = RGBColor(0x21, 0x26, 0x2d)   # #21262d
GREEN    = RGBColor(0x3f, 0xb9, 0x50)   # #3fb950
BLUE     = RGBColor(0x58, 0xa6, 0xff)   # #58a6ff
YELLOW   = RGBColor(0xe3, 0xb3, 0x41)   # #e3b341
RED      = RGBColor(0xf8, 0x51, 0x49)   # #f85149
PURPLE   = RGBColor(0xa7, 0x8b, 0xfa)   # #a78bfa
WHITE    = RGBColor(0xf0, 0xf6, 0xfc)   # #f0f6fc
GREY     = RGBColor(0x8b, 0x94, 0x9e)   # #8b949e
DIMGREY  = RGBColor(0x48, 0x4f, 0x58)   # #484f58

W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def box(slide, left, top, width, height, color=None, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    if color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    else:
        shape.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "JetBrains Mono" if size <= 14 else "Calibri"
    return txb


def accent_bar(slide, color=BLUE):
    """Top accent line."""
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def slide_number(slide, n, total=10):
    txt(slide, f"{n:02d} / {total:02d}", 12.3, 7.1, 1.0, 0.35,
        size=8, color=DIMGREY, align=PP_ALIGN.RIGHT)


def tag(slide, label, left, top, color=BLUE):
    b = box(slide, left, top, len(label)*0.085 + 0.2, 0.28, color=color)
    t = slide.shapes.add_textbox(Inches(left+0.07), Inches(top+0.04),
                                  Inches(len(label)*0.085+0.06), Inches(0.22))
    p = t.text_frame.paragraphs[0]
    r = p.add_run(); r.text = label
    r.font.size = Pt(7.5); r.font.bold = True
    r.font.color.rgb = BG; r.font.name = "Calibri"


# ════════════════════════════════════════════════════════════════════
#  SLIDES
# ════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    sl = blank_slide(prs); bg(sl)
    accent_bar(sl, BLUE)

    # Decorative right panel
    b = box(sl, 8.8, 0.06, 4.53, 7.44, PANEL)
    b.line.color.rgb = BORDER; b.line.width = Pt(0.5)

    # Grid lines decoration
    for i in range(6):
        g = sl.shapes.add_shape(1, Inches(8.8), Inches(0.06 + i*1.24),
                                 Inches(4.53), Inches(0.02))
        g.fill.solid(); g.fill.fore_color.rgb = BORDER
        g.line.fill.background()

    # Terminal mockup text
    lines = [
        ("$ crypto-terminal --start", GREEN),
        ("", WHITE),
        ("[Ollama]  Mistral 7B connected", BLUE),
        ("[XGBoost] Models loaded (5 coins)", BLUE),
        ("[Pipeline] LIVE — 44,185 candles", GREEN),
        ("[Signal]  BTC  HOLD  conf=0.62", YELLOW),
        ("[Signal]  ETH  HOLD  conf=0.58", YELLOW),
        ("[Whale]   0.01 BTC whale alert", RED),
        ("", WHITE),
        ("Dashboard  →  localhost:8501", GREY),
    ]
    for i, (line, c) in enumerate(lines):
        txt(sl, line, 9.0, 0.5 + i*0.65, 4.1, 0.55, size=9, color=c)

    # Event badge
    tag(sl, "NMIMS INNOVATHON 2026", 0.5, 0.3, YELLOW)
    tag(sl, "CHALLENGE 2 OF 5", 3.1, 0.3, PURPLE)

    txt(sl, "CRYPTO TRADING", 0.5, 1.1, 8.0, 1.1, size=48, bold=True, color=WHITE)
    txt(sl, "INTELLIGENCE TERMINAL", 0.5, 1.95, 8.0, 1.0, size=34, bold=True, color=BLUE)

    txt(sl, "Self-Hosted LLM · Price Prediction · On-Chain Analytics · Signal Engine",
        0.5, 3.15, 8.0, 0.5, size=13, color=GREY)

    # Stat pills
    stats = [("44K+", "Price Candles", GREEN), ("5", "Coins Tracked", BLUE),
             ("5,577", "Signals Generated", YELLOW), ("₹0", "Total Cost", PURPLE)]
    for i, (val, lbl, c) in enumerate(stats):
        bx = box(sl, 0.5 + i*2.0, 4.0, 1.8, 0.9, PANEL)
        bx.line.color.rgb = c; bx.line.width = Pt(1)
        txt(sl, val, 0.55 + i*2.0, 4.05, 1.7, 0.45, size=20, bold=True, color=c, align=PP_ALIGN.CENTER)
        txt(sl, lbl, 0.55 + i*2.0, 4.45, 1.7, 0.35, size=8, color=GREY, align=PP_ALIGN.CENTER)

    txt(sl, "AI & Data Intelligence Track  ·  HIGH Difficulty  ·  24-Hour Build",
        0.5, 5.2, 8.0, 0.4, size=10, color=DIMGREY)

    slide_number(sl, 1)
    return sl


def slide_02_problem(prs):
    sl = blank_slide(prs); bg(sl)
    accent_bar(sl, RED)
    tag(sl, "THE PROBLEM", 0.5, 0.15, RED)
    txt(sl, "Retail traders have no affordable intelligence tools", 0.5, 0.55, 12.5, 0.7,
        size=28, bold=True, color=WHITE)

    problems = [
        ("₹5L/month", "Santiment costs ₹5 lakh per month — inaccessible to retail traders"),
        ("Manual work", "Traders spend hours reading Reddit, Telegram, news manually"),
        ("No integration", "Existing bots use technical indicators only — ignore sentiment & on-chain"),
        ("No local LLM", "Cloud APIs are expensive and leak sensitive trading strategies"),
    ]
    for i, (title, desc) in enumerate(problems):
        row = 1.5 + i * 1.3
        b = box(sl, 0.5, row, 12.3, 1.1, PANEL)
        b.line.color.rgb = RED; b.line.width = Pt(0.8)
        txt(sl, title, 0.65, row+0.08, 2.0, 0.4, size=11, bold=True, color=RED)
        txt(sl, desc,  0.65, row+0.5,  11.9, 0.5, size=11, color=GREY)

    txt(sl, "Professional platforms: Santiment ₹5L/mo · LumiBot ₹2L/mo · CryptoQuant ₹1.5L/mo",
        0.5, 6.9, 12.3, 0.4, size=9, color=DIMGREY)
    slide_number(sl, 2)


def slide_03_solution(prs):
    sl = blank_slide(prs); bg(sl)
    accent_bar(sl, GREEN)
    tag(sl, "OUR SOLUTION", 0.5, 0.15, GREEN)
    txt(sl, "A fully self-hosted crypto intelligence terminal at ₹0", 0.5, 0.55, 12.5, 0.65,
        size=26, bold=True, color=WHITE)

    pillars = [
        (GREEN,  "LLM Sentiment",    "Mistral 7B via Ollama\nFinBERT fallback\nBULLISH / BEARISH / NEUTRAL"),
        (BLUE,   "Price Prediction", "Facebook Prophet\nXGBoost (27 features)\n1h, 4h, 24h horizons"),
        (YELLOW, "On-Chain Data",    "Whale Alert API\n20+ chains tracked\nAccumulation vs Distribution"),
        (PURPLE, "Signal Engine",    "3-factor rule engine\nBUY / SELL / HOLD\nConfidence + Reasoning"),
    ]
    for i, (c, title, body) in enumerate(pillars):
        col = 0.4 + i * 3.22
        b = box(sl, col, 1.5, 3.0, 3.5, PANEL)
        b.line.color.rgb = c; b.line.width = Pt(1.5)
        # top color bar
        cb = sl.shapes.add_shape(1, Inches(col), Inches(1.5), Inches(3.0), Inches(0.08))
        cb.fill.solid(); cb.fill.fore_color.rgb = c; cb.line.fill.background()
        txt(sl, title, col+0.1, 1.65, 2.8, 0.45, size=13, bold=True, color=c)
        txt(sl, body,  col+0.1, 2.2,  2.8, 2.5,  size=10, color=GREY)

    txt(sl, "Total infrastructure cost: ₹0  ·  Runs on 8GB RAM  ·  No cloud APIs required",
        0.5, 5.3, 12.3, 0.4, size=10, color=GREEN, align=PP_ALIGN.CENTER)

    # vs competitors
    b = box(sl, 0.5, 5.8, 12.3, 1.0, PANEL)
    b.line.color.rgb = BORDER
    comps = [("Santiment", "₹5L/mo", RED), ("LumiBot", "₹2L/mo", RED),
             ("Our Terminal", "₹0", GREEN)]
    for i, (name, price, c) in enumerate(comps):
        x = 1.5 + i * 3.8
        txt(sl, name,  x, 5.85, 3.0, 0.35, size=10, color=GREY,  align=PP_ALIGN.CENTER)
        txt(sl, price, x, 6.2,  3.0, 0.45, size=18, bold=True, color=c, align=PP_ALIGN.CENTER)

    slide_number(sl, 3)


def slide_04_architecture(prs):
    sl = blank_slide(prs); bg(sl)
    accent_bar(sl, PURPLE)
    tag(sl, "ARCHITECTURE", 0.5, 0.15, PURPLE)
    txt(sl, "End-to-end pipeline: Data → Intelligence → Signal → Dashboard",
        0.5, 0.55, 12.5, 0.55, size=22, bold=True, color=WHITE)

    # Layer boxes
    layers = [
        (BLUE,   "DATA SOURCES",
         "Telegram Channels  ·  NewsAPI  ·  Whale Alert API  ·  Binance REST (15-min OHLCV)"),
        (YELLOW, "INGESTION PIPELINE",
         "Async collectors  ·  Exponential backoff retry  ·  SQLAlchemy ORM  ·  SQLite (7 tables)"),
        (GREEN,  "INTELLIGENCE LAYER",
         "Ollama Mistral 7B sentiment  ·  Prophet time-series  ·  XGBoost (27 features)  ·  FinBERT fallback"),
        (PURPLE, "SIGNAL ENGINE",
         "3-factor rule: Sentiment + ML Prediction + Whale Flow  →  BUY / SELL / HOLD + confidence"),
        (RED,    "OUTPUT",
         "Streamlit dashboard (6 tabs)  ·  Backtesting engine  ·  Paper trading simulator"),
    ]
    for i, (c, title, desc) in enumerate(layers):
        row = 1.3 + i * 1.1
        b = box(sl, 0.5, row, 12.3, 0.95, PANEL)
        b.line.color.rgb = c; b.line.width = Pt(1)
        lb = sl.shapes.add_shape(1, Inches(0.5), Inches(row), Inches(0.08), Inches(0.95))
        lb.fill.solid(); lb.fill.fore_color.rgb = c; lb.line.fill.background()
        txt(sl, title, 0.75, row+0.07, 2.5, 0.35, size=8, bold=True, color=c)
        txt(sl, desc,  0.75, row+0.45, 11.8, 0.4, size=10, color=GREY)

        if i < 4:
            arr = sl.shapes.add_shape(1, Inches(6.4), Inches(row+0.95),
                                       Inches(0.5), Inches(0.15))
            arr.fill.solid(); arr.fill.fore_color.rgb = c; arr.line.fill.background()

    slide_number(sl, 4)


def slide_05_sentiment(prs):
    sl = blank_slide(prs); bg(sl)
    accent_bar(sl, YELLOW)
    tag(sl, "SENTIMENT ANALYSIS", 0.5, 0.15, YELLOW)
    txt(sl, "Self-hosted LLM — Mistral 7B via Ollama", 0.5, 0.55, 9.0, 0.6,
        size=26, bold=True, color=WHITE)

    # Left: flow
    steps = [
        (BLUE,   "1. Collect", "Telegram channels + NewsAPI\nReal-time crypto content"),
        (YELLOW, "2. Classify", "Mistral 7B prompt:\nBULLISH / BEARISH / NEUTRAL\nScore: 0.0 – 1.0"),
        (GREEN,  "3. Fallback", "FinBERT (CPU-based)\nActivates if Ollama unavailable"),
        (PURPLE, "4. Aggregate", "Rolling avg per coin\nFear & Greed Index overlay"),
    ]
    for i, (c, title, body) in enumerate(steps):
        row = 1.4 + i * 1.35
        b = box(sl, 0.4, row, 5.5, 1.15, PANEL)
        b.line.color.rgb = c; b.line.width = Pt(1)
        txt(sl, title, 0.6, row+0.1, 5.2, 0.35, size=11, bold=True, color=c)
        txt(sl, body,  0.6, row+0.5, 5.2, 0.55, size=9,  color=GREY)

    # Right: thresholds + live data
    b = box(sl, 6.4, 1.4, 6.5, 4.1, PANEL)
    b.line.color.rgb = BORDER
    txt(sl, "SIGNAL THRESHOLDS", 6.6, 1.5, 6.1, 0.35, size=9, bold=True, color=DIMGREY)

    thresholds = [
        ("> 0.70", "BULLISH", GREEN,  "Strong buy signal"),
        ("0.30–0.70", "NEUTRAL", YELLOW, "No directional signal"),
        ("< 0.30", "BEARISH", RED,   "Strong sell signal"),
    ]
    for i, (rng, label, c, note) in enumerate(thresholds):
        row = 2.0 + i * 1.0
        txt(sl, rng,   6.6, row,      2.0, 0.35, size=13, bold=True, color=c)
        txt(sl, label, 8.8, row,      2.0, 0.35, size=13, bold=True, color=c)
        txt(sl, note,  6.6, row+0.38, 5.5, 0.3,  size=9,  color=GREY)

    txt(sl, "Live reading: 0.12 — EXTREME FEAR (Fear & Greed: 12/100)",
        6.6, 5.1, 6.1, 0.35, size=9, color=RED)
    txt(sl, "2,333 sentiment snapshots collected",
        6.6, 5.5, 6.1, 0.3, size=9, color=GREY)

    slide_number(sl, 5)


def slide_06_ml(prs):
    sl = blank_slide(prs); bg(sl)
    accent_bar(sl, BLUE)
    tag(sl, "PRICE PREDICTION", 0.5, 0.15, BLUE)
    txt(sl, "Prophet + XGBoost — 27 engineered features", 0.5, 0.55, 12.5, 0.6,
        size=26, bold=True, color=WHITE)

    # Accuracy table
    txt(sl, "XGBoost Directional Accuracy (test set)", 0.5, 1.35, 6.5, 0.4,
        size=11, bold=True, color=DIMGREY)

    results = [
        ("BTC", 61.1, BLUE),
        ("ETH", 56.0, BLUE),
        ("SOL", 67.3, GREEN),
        ("XRP", 65.0, GREEN),
        ("DOGE", 61.4, BLUE),
    ]
    for i, (coin, acc, c) in enumerate(results):
        row = 1.85 + i * 0.82
        b = box(sl, 0.5, row, 6.1, 0.7, PANEL)
        b.line.color.rgb = BORDER
        txt(sl, coin, 0.65, row+0.17, 1.0, 0.35, size=12, bold=True, color=WHITE)
        # bar
        bar_w = (acc - 50) / 50 * 4.2
        br = sl.shapes.add_shape(1, Inches(1.7), Inches(row+0.22),
                                  Inches(bar_w), Inches(0.28))
        br.fill.solid(); br.fill.fore_color.rgb = c; br.line.fill.background()
        txt(sl, f"{acc:.1f}%", 5.9, row+0.17, 0.75, 0.35, size=12, bold=True, color=c, align=PP_ALIGN.RIGHT)

    # 55% target line label
    txt(sl, "▲ 55% target", 3.5, 6.1, 2.0, 0.3, size=8, color=YELLOW)

    # Right: feature groups
    b = box(sl, 7.1, 1.35, 5.8, 5.3, PANEL)
    b.line.color.rgb = BORDER
    txt(sl, "FEATURE GROUPS", 7.3, 1.45, 5.4, 0.35, size=9, bold=True, color=DIMGREY)

    fgroups = [
        ("Technical Indicators", "22", BLUE,
         "SMA, EMA, MACD, RSI, Bollinger Bands,\nATR, volume ratio, returns, volatility"),
        ("Sentiment Signals", "3", GREEN,
         "sentiment_avg, sentiment_momentum,\nsentiment_volume"),
        ("On-Chain / Whale", "2", YELLOW,
         "whale_net_flow, whale_tx_count"),
    ]
    for i, (name, count, c, detail) in enumerate(fgroups):
        row = 2.0 + i * 1.5
        txt(sl, name,   7.3, row,      3.5, 0.35, size=10, bold=True, color=c)
        txt(sl, count,  11.8, row,     0.9, 0.35, size=14, bold=True, color=c, align=PP_ALIGN.RIGHT)
        txt(sl, detail, 7.3, row+0.4,  5.4, 0.7,  size=8,  color=GREY)

    txt(sl, "Total: 27 Features", 7.3, 6.1, 5.4, 0.35, size=11, bold=True, color=PURPLE)

    slide_number(sl, 6)


def slide_07_signal(prs):
    sl = blank_slide(prs); bg(sl)
    accent_bar(sl, GREEN)
    tag(sl, "SIGNAL ENGINE", 0.5, 0.15, GREEN)
    txt(sl, "Three independent factors — two must agree", 0.5, 0.55, 12.5, 0.6,
        size=26, bold=True, color=WHITE)

    # 3 factor boxes
    factors = [
        (GREEN,  "SENTIMENT",   "Mistral 7B score\n> 0.70 = Bullish\n< 0.30 = Bearish"),
        (BLUE,   "ML PREDICTION", "XGBoost / Prophet\nUP / DOWN / SIDEWAYS\n+ Confidence %"),
        (YELLOW, "WHALE FLOW",  "Whale Alert API\nAccumulation vs\nDistribution"),
    ]
    for i, (c, title, body) in enumerate(factors):
        col = 0.5 + i * 3.0
        b = box(sl, col, 1.4, 2.7, 2.5, PANEL)
        b.line.color.rgb = c; b.line.width = Pt(2)
        txt(sl, title, col+0.1, 1.5,  2.5, 0.45, size=12, bold=True, color=c)
        txt(sl, body,  col+0.1, 2.05, 2.5, 1.7,  size=10, color=GREY)

    # Arrow + combine box
    txt(sl, "+", 3.3, 2.25, 0.4, 0.6, size=24, bold=True, color=DIMGREY, align=PP_ALIGN.CENTER)
    txt(sl, "+", 6.3, 2.25, 0.4, 0.6, size=24, bold=True, color=DIMGREY, align=PP_ALIGN.CENTER)

    txt(sl, "▼  2 of 3 factors must align", 4.5, 4.05, 4.5, 0.4,
        size=10, color=DIMGREY, align=PP_ALIGN.CENTER)

    signals = [
        (GREEN, "BUY",  "Sentiment > 0.70  AND  Prediction = UP  AND  Whales accumulating"),
        (RED,   "SELL", "Sentiment < 0.30  AND  Prediction = DOWN  AND  Whales distributing"),
        (GREY,  "HOLD", "Any other combination — no strong consensus"),
    ]
    for i, (c, sig, rule) in enumerate(signals):
        row = 4.6 + i * 0.85
        b = box(sl, 0.5, row, 12.3, 0.72, PANEL)
        b.line.color.rgb = c; b.line.width = Pt(1)
        sb = sl.shapes.add_shape(1, Inches(0.5), Inches(row), Inches(0.08), Inches(0.72))
        sb.fill.solid(); sb.fill.fore_color.rgb = c; sb.line.fill.background()
        txt(sl, sig,  0.75, row+0.17, 1.0, 0.38, size=13, bold=True, color=c)
        txt(sl, rule, 1.85, row+0.17, 10.8, 0.38, size=10, color=GREY)

    slide_number(sl, 7)


def slide_08_backtest(prs):
    sl = blank_slide(prs); bg(sl)
    accent_bar(sl, YELLOW)
    tag(sl, "BACKTESTING", 0.5, 0.15, YELLOW)
    txt(sl, "Realistic signal replay with cost modeling", 0.5, 0.55, 12.5, 0.6,
        size=26, bold=True, color=WHITE)

    # Assumptions
    b = box(sl, 0.5, 1.35, 5.8, 2.2, PANEL)
    b.line.color.rgb = BORDER
    txt(sl, "SIMULATION ASSUMPTIONS", 0.7, 1.45, 5.4, 0.35, size=9, bold=True, color=DIMGREY)
    assumptions = [
        ("Slippage",    "0.1% per side"),
        ("Commission",  "0.1% per trade"),
        ("Position Size", "Confidence-weighted"),
        ("Data Window", "90 days OHLCV"),
    ]
    for i, (k, v) in enumerate(assumptions):
        row = 1.9 + i * 0.45
        txt(sl, k, 0.7, row, 2.5, 0.35, size=10, color=GREY)
        txt(sl, v, 3.5, row, 2.7, 0.35, size=10, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    # Metrics
    b2 = box(sl, 6.7, 1.35, 6.1, 2.2, PANEL)
    b2.line.color.rgb = BORDER
    txt(sl, "METRICS COMPUTED", 6.9, 1.45, 5.7, 0.35, size=9, bold=True, color=DIMGREY)
    metrics = ["Total Return", "CAGR", "Sharpe Ratio", "Sortino Ratio", "Max Drawdown", "Win Rate"]
    for i, m in enumerate(metrics):
        col = 6.9 + (i % 2) * 2.9
        row = 1.9 + (i // 2) * 0.45
        txt(sl, f"• {m}", col, row, 2.8, 0.35, size=10, color=GREY)

    # Pipeline
    txt(sl, "HOW IT WORKS", 0.5, 3.75, 12.0, 0.35, size=9, bold=True, color=DIMGREY)
    steps = [
        (YELLOW, "Load Signals", "Query historical BUY/SELL\nsignals from DB"),
        (BLUE,   "Replay Prices", "Match each signal to\nOHLCV price data"),
        (GREEN,  "Simulate Fill", "Apply slippage +\ncommission at entry/exit"),
        (PURPLE, "Compute Risk", "Sharpe, Sortino,\ndrawdown, CAGR"),
    ]
    for i, (c, title, body) in enumerate(steps):
        col = 0.5 + i * 3.2
        b = box(sl, col, 4.15, 2.9, 2.0, PANEL)
        b.line.color.rgb = c; b.line.width = Pt(1)
        txt(sl, title, col+0.1, 4.25, 2.7, 0.4,  size=11, bold=True, color=c)
        txt(sl, body,  col+0.1, 4.75, 2.7, 0.9,  size=9,  color=GREY)

    slide_number(sl, 8)


def slide_09_dashboard(prs):
    sl = blank_slide(prs); bg(sl)
    accent_bar(sl, PURPLE)
    tag(sl, "DASHBOARD", 0.5, 0.15, PURPLE)
    txt(sl, "6-tab Streamlit terminal — live at localhost:8501", 0.5, 0.55, 12.5, 0.6,
        size=26, bold=True, color=WHITE)

    tabs = [
        (BLUE,   "01  Market Overview",      "Candlestick charts, OHLCV, RSI,\nMarket Regime panel, signal cards"),
        (YELLOW, "02  Sentiment Intel",       "Heatmap, trend chart, Fear & Greed,\nRecent news with LLM scores"),
        (GREEN,  "03  Predictions",           "ML Intelligence panel, Prophet &\nXGBoost confidence, signal log"),
        (RED,    "04  Whale Tracker",         "Live whale transactions,\nAccumulation vs Distribution"),
        (PURPLE, "05  Backtesting",           "Equity curve, Sharpe ratio,\nMax drawdown, trade log"),
        (WHITE,  "06  Paper Trading",         "Execute simulated trades,\nAuto-close at 4h, P&L tracking"),
    ]
    for i, (c, title, body) in enumerate(tabs):
        col = 0.4 + (i % 3) * 4.3
        row = 1.5 + (i // 3) * 2.3
        b = box(sl, col, row, 4.0, 2.0, PANEL)
        b.line.color.rgb = c; b.line.width = Pt(1.5)
        cb = sl.shapes.add_shape(1, Inches(col), Inches(row), Inches(4.0), Inches(0.07))
        cb.fill.solid(); cb.fill.fore_color.rgb = c; cb.line.fill.background()
        txt(sl, title, col+0.12, row+0.15, 3.76, 0.4, size=11, bold=True, color=c)
        txt(sl, body,  col+0.12, row+0.65, 3.76, 1.1, size=9,  color=GREY)

    txt(sl, "Auto-refresh · REFRESH NOW button · Pipeline status indicator · IST timestamps",
        0.5, 7.05, 12.3, 0.35, size=9, color=DIMGREY, align=PP_ALIGN.CENTER)

    slide_number(sl, 9)


def slide_10_stack(prs):
    sl = blank_slide(prs); bg(sl)
    accent_bar(sl, GREEN)
    tag(sl, "TECH STACK & SUMMARY", 0.5, 0.15, GREEN)
    txt(sl, "100% open-source · ₹0 cost · runs on 8GB RAM", 0.5, 0.55, 10.0, 0.6,
        size=24, bold=True, color=WHITE)

    stack = [
        ("LLM Sentiment",    "Ollama + Mistral 7B Q4 + FinBERT fallback",       YELLOW),
        ("Price Prediction", "Facebook Prophet + XGBoost",                        BLUE),
        ("Feature Eng.",     "27 features: technical, sentiment, on-chain",       BLUE),
        ("Data Collection",  "Telethon, NewsAPI, Whale Alert, Binance REST",      GREEN),
        ("Signal Engine",    "Custom 3-factor rule engine",                        GREEN),
        ("Backtesting",      "Slippage + commission + Sharpe/Sortino/drawdown",   YELLOW),
        ("Paper Trading",    "Local simulation, live prices, auto-close 4h",      PURPLE),
        ("Dashboard",        "Streamlit + Plotly, 6 tabs",                         PURPLE),
        ("Database",         "SQLAlchemy ORM + SQLite (7 tables, indexed)",        GREY),
        ("Resilience",       "Exponential backoff retry on all API collectors",    GREY),
    ]
    for i, (comp, tech, c) in enumerate(stack):
        col = 0.4 + (i % 2) * 6.4
        row = 1.4 + (i // 2) * 0.72
        b = box(sl, col, row, 6.1, 0.6, PANEL)
        b.line.color.rgb = BORDER
        lb = sl.shapes.add_shape(1, Inches(col), Inches(row), Inches(0.06), Inches(0.6))
        lb.fill.solid(); lb.fill.fore_color.rgb = c; lb.line.fill.background()
        txt(sl, comp, col+0.15, row+0.1, 1.8, 0.35, size=9,  bold=True, color=c)
        txt(sl, tech, col+2.05, row+0.1, 3.9, 0.35, size=9,  color=GREY)

    # Final tagline
    b = box(sl, 0.5, 7.0, 12.3, 0.38, GREEN)
    b.line.fill.background()
    txt(sl, "Replicates Santiment / Nansen capabilities — at zero cost, running entirely on local hardware",
        0.6, 7.04, 12.1, 0.3, size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)

    slide_number(sl, 10)


# ── Build ────────────────────────────────────────────────────────────────────
prs = new_prs()
slide_01_title(prs)
slide_02_problem(prs)
slide_03_solution(prs)
slide_04_architecture(prs)
slide_05_sentiment(prs)
slide_06_ml(prs)
slide_07_signal(prs)
slide_08_backtest(prs)
slide_09_dashboard(prs)
slide_10_stack(prs)

out = "Crypto_Terminal_Innovathon2026.pptx"
prs.save(out)
print(f"Saved: {out}")
